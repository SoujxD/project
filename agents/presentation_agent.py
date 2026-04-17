"""Presentation generator agent for dataset-adaptive insight decks."""

from __future__ import annotations

import math
import shutil
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from utils.dataset_adapter import load_analysis_dataset


@dataclass(slots=True)
class SlideContent:
    """Structured content for a slide."""

    title: str
    bullets: list[str]
    speaker_notes: str
    chart_key: str | None = None


class PresentationGeneratorAgent:
    """Generates dataset-focused charts and a reusable PowerPoint deck."""

    CHART_DPI = 180

    def __init__(self, output_dir: str | Path) -> None:
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.chart_dir = self.output_dir / "charts"
        self.chart_dir.mkdir(parents=True, exist_ok=True)
        self.docs_asset_dir = self.output_dir.parent / "docs" / "assets"
        self.docs_asset_dir.mkdir(parents=True, exist_ok=True)

    def _normalized(self, name: str) -> str:
        return "".join(ch.lower() for ch in str(name) if ch.isalnum())

    def _find_column(self, dataset: pd.DataFrame, keywords: list[str], exclude: set[str] | None = None) -> str | None:
        exclude = exclude or set()
        normalized_map = {column: self._normalized(column) for column in dataset.columns}
        for keyword in keywords:
            probe = self._normalized(keyword)
            for column, normalized in normalized_map.items():
                if column in exclude:
                    continue
                if probe in normalized:
                    return column
        return None

    def _target_series(self, series: pd.Series) -> pd.Series | None:
        if pd.api.types.is_bool_dtype(series):
            return series.astype(int)
        if pd.api.types.is_numeric_dtype(series):
            cleaned = pd.to_numeric(series, errors="coerce").dropna()
            unique_values = set(cleaned.unique().tolist())
            if unique_values and unique_values.issubset({0, 1}):
                return pd.to_numeric(series, errors="coerce").fillna(0).astype(int)
        non_null = series.dropna()
        lowered = non_null.astype("string").str.strip().str.lower()
        mapping = {
            "1": 1,
            "0": 0,
            "true": 1,
            "false": 0,
            "yes": 1,
            "no": 0,
            "y": 1,
            "n": 0,
            "purchase": 1,
            "no_purchase": 0,
            "purchased": 1,
            "not_purchased": 0,
            "converted": 1,
            "not_converted": 0,
            "retained": 1,
            "not_retained": 0,
            "churned": 1,
            "not_churned": 0,
            "approved": 1,
            "rejected": 0,
            "success": 1,
            "failed": 0,
        }
        unique_values = set(lowered.dropna().unique().tolist())
        if unique_values and unique_values.issubset(set(mapping)):
            normalized = series.astype("string").str.strip().str.lower()
            return normalized.map(mapping).fillna(0).astype(int)
        return None

    def _is_identifierish(self, column: str, series: pd.Series) -> bool:
        normalized = self._normalized(column)
        if any(token in normalized for token in ["id", "uuid", "identifier"]):
            return True
        if not pd.api.types.is_numeric_dtype(series):
            return False
        numeric = pd.to_numeric(series, errors="coerce").dropna()
        if numeric.empty or numeric.nunique(dropna=True) != len(numeric):
            return False
        ordered = numeric.sort_values().reset_index(drop=True)
        steps = ordered.diff().dropna()
        return not steps.empty and steps.nunique() == 1

    def _resolve_target_column(self, dataset: pd.DataFrame) -> str | None:
        keyword_matches: list[str] = []
        normalized_map = {column: self._normalized(column) for column in dataset.columns}
        keywords = [
            "revenue",
            "target",
            "label",
            "outcome",
            "converted",
            "conversion",
            "purchased",
            "purchase",
            "churn",
            "retain",
            "default",
            "fraud",
            "approved",
            "success",
        ]
        for keyword in keywords:
            probe = self._normalized(keyword)
            keyword_matches.extend(
                column
                for column, normalized in normalized_map.items()
                if probe in normalized and column not in keyword_matches
            )
        for column in keyword_matches:
            if self._target_series(dataset[column]) is not None:
                return column
        for column in dataset.columns:
            if self._target_series(dataset[column]) is not None:
                return column
        return None

    def _time_like_column(self, dataset: pd.DataFrame, exclude: set[str]) -> str | None:
        keyword_match = self._find_column(
            dataset,
            ["timestamp", "datetime", "date", "month", "week", "day", "year", "period", "quarter", "time"],
            exclude=exclude,
        )
        if keyword_match:
            return keyword_match

        for column in dataset.columns:
            if column in exclude:
                continue
            series = dataset[column]
            if pd.api.types.is_datetime64_any_dtype(series):
                return column
            if not (pd.api.types.is_object_dtype(series) or pd.api.types.is_string_dtype(series)):
                continue
            sample = series.dropna().astype("string").head(200)
            if sample.empty:
                continue
            parsed = pd.to_datetime(sample, errors="coerce", format="mixed")
            success_rate = float(parsed.notna().mean())
            if success_rate >= 0.7 and sample.nunique() >= 3:
                return column
        return None

    def _categorical_candidates(self, dataset: pd.DataFrame, exclude: set[str]) -> list[str]:
        candidates: list[str] = []
        for column in dataset.columns:
            if column in exclude:
                continue
            series = dataset[column]
            nunique = int(series.nunique(dropna=True))
            if nunique <= 1:
                continue
            if pd.api.types.is_bool_dtype(series) or pd.api.types.is_object_dtype(series) or pd.api.types.is_string_dtype(series):
                candidates.append(column)
                continue
            if self._is_identifierish(column, series):
                continue
            if pd.api.types.is_numeric_dtype(series) and 2 < nunique <= min(12, max(4, int(len(dataset) * 0.25))):
                candidates.append(column)
        return candidates

    def _numeric_metric_candidates(self, dataset: pd.DataFrame, exclude: set[str]) -> list[str]:
        scored: list[tuple[float, str]] = []
        for column in dataset.select_dtypes(include=["number", "bool"]).columns:
            if column in exclude:
                continue
            if self._is_identifierish(column, dataset[column]):
                continue
            series = pd.to_numeric(dataset[column], errors="coerce")
            if series.dropna().empty:
                continue
            unique_values = set(series.dropna().unique().tolist())
            if unique_values and unique_values.issubset({0, 1}):
                continue
            nunique = series.nunique(dropna=True)
            if nunique <= 3:
                continue
            normalized = self._normalized(column)
            score = float(min(nunique, 50))
            if any(token in normalized for token in ["value", "rate", "score", "amount", "price", "duration", "page", "bounce", "exit"]):
                score += 25
            if any(token in normalized for token in ["type", "region", "category", "segment"]):
                score -= 10
            scored.append((score, column))
        return [column for _, column in sorted(scored, reverse=True)]

    def infer_schema(self, dataset: pd.DataFrame) -> dict[str, Any]:
        """Infer broadly useful roles from a flexible tabular dataset."""
        target_col = self._resolve_target_column(dataset)
        time_col = self._time_like_column(dataset, exclude={target_col} - {None})

        exclude = {target_col, time_col} - {None}
        categorical_candidates = self._categorical_candidates(dataset, exclude)
        numeric_candidates = self._numeric_metric_candidates(dataset, exclude)

        primary_dimension = categorical_candidates[0] if categorical_candidates else None
        secondary_dimension = categorical_candidates[1] if len(categorical_candidates) > 1 else None
        primary_metric = numeric_candidates[0] if numeric_candidates else None
        secondary_metric = numeric_candidates[1] if len(numeric_candidates) > 1 else None

        return {
            "target": target_col,
            "time": time_col,
            "primary_dimension": primary_dimension,
            "secondary_dimension": secondary_dimension,
            "primary_metric": primary_metric,
            "secondary_metric": secondary_metric,
            "categorical_candidates": categorical_candidates[:6],
            "numeric_candidates": numeric_candidates[:6],
        }

    def _target_metric(self, dataset: pd.DataFrame, schema: dict[str, Any]) -> pd.Series | None:
        target_col = schema.get("target")
        if not target_col or target_col not in dataset.columns:
            return None
        return self._target_series(dataset[target_col])

    def _series_by_dimension(
        self,
        dataset: pd.DataFrame,
        schema: dict[str, Any],
        dimension: str | None,
        top_n: int = 8,
    ) -> pd.Series | None:
        if not dimension or dimension not in dataset.columns:
            return None
        working = dataset.copy()
        metric = self._target_metric(working, schema)
        if metric is not None:
            working["_metric_target"] = metric
            return working.groupby(dimension, dropna=False)["_metric_target"].mean().sort_values(ascending=False).head(top_n)
        return working[dimension].astype(str).value_counts(dropna=False).head(top_n).sort_values(ascending=False)

    def _time_series(self, dataset: pd.DataFrame, schema: dict[str, Any], top_n: int = 12) -> pd.Series | None:
        time_col = schema.get("time")
        if not time_col or time_col not in dataset.columns:
            return None
        metric = self._target_metric(dataset, schema)
        labels = dataset[time_col]
        if pd.api.types.is_datetime64_any_dtype(labels):
            labels = labels.dt.strftime("%Y-%m-%d")
        else:
            converted = pd.to_datetime(labels, errors="coerce", format="mixed")
            if converted.notna().mean() >= 0.6:
                labels = converted.dt.strftime("%Y-%m-%d").fillna(dataset[time_col].astype(str))
            else:
                labels = dataset[time_col].astype(str)
        working = pd.DataFrame({"time": labels})
        if metric is not None:
            working["metric"] = metric
            series = working.groupby("time", dropna=False)["metric"].mean()
        else:
            series = working["time"].value_counts(dropna=False)
        return series.sort_index().tail(top_n)

    def _bucket_series(self, dataset: pd.DataFrame, schema: dict[str, Any]) -> tuple[pd.Series | None, str]:
        numeric_col = schema.get("primary_metric") or schema.get("secondary_metric")
        if not numeric_col or numeric_col not in dataset.columns:
            return None, ""
        numeric_series = pd.to_numeric(dataset[numeric_col], errors="coerce")
        valid = numeric_series.dropna()
        if valid.empty:
            return None, ""
        quantiles = np.unique(np.quantile(valid, [0, 0.25, 0.5, 0.75, 1.0]))
        bins = quantiles if len(quantiles) >= 3 else np.linspace(float(valid.min()), float(valid.max()) + 1e-6, 5)
        labels = [f"Q{i + 1}" for i in range(len(bins) - 1)]
        bucketed = pd.cut(numeric_series, bins=bins, labels=labels, include_lowest=True, duplicates="drop")
        metric = self._target_metric(dataset, schema)
        if metric is not None:
            working = pd.DataFrame({"bucket": bucketed, "metric": metric})
            return working.groupby("bucket", observed=False)["metric"].mean(), numeric_col
        return bucketed.value_counts().sort_index(), numeric_col

    def _count_missing(self, dataset: pd.DataFrame) -> pd.Series:
        return (dataset.isna().mean() * 100).sort_values(ascending=False)

    def _plot_bar(self, series: pd.Series, title: str, path: Path, horizontal: bool = False, ylabel: str = "Value") -> None:
        fig, ax = plt.subplots(figsize=(8, 4.4))
        color = "#1E5F74"
        if horizontal:
            series.sort_values().plot(kind="barh", ax=ax, color=color)
            ax.set_xlabel(ylabel)
            ax.set_ylabel("")
        else:
            series.plot(kind="bar", ax=ax, color=color)
            ax.set_xlabel("")
            ax.set_ylabel(ylabel)
            plt.setp(ax.get_xticklabels(), rotation=24, ha="right")
        ax.set_title(title, fontsize=13)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        fig.savefig(path, dpi=self.CHART_DPI)
        plt.close(fig)

    def _plot_line(self, series: pd.Series, title: str, path: Path, ylabel: str = "Value") -> None:
        cleaned = pd.to_numeric(series, errors="coerce").dropna()
        if cleaned.empty:
            return
        fig, ax = plt.subplots(figsize=(8, 4.4))
        x = np.arange(len(cleaned))
        ax.plot(x, cleaned.values.astype(float), color="#1E5F74", linewidth=2.5, marker="o")
        ax.fill_between(x, cleaned.values.astype(float), color="#9ED8DB", alpha=0.5)
        ax.set_xticks(x)
        ax.set_xticklabels([str(value) for value in cleaned.index], rotation=22, ha="right")
        ax.set_title(title, fontsize=13)
        ax.set_xlabel("")
        ax.set_ylabel(ylabel)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        fig.savefig(path, dpi=self.CHART_DPI)
        plt.close(fig)

    def _plot_hist(self, series: pd.Series, title: str, path: Path) -> None:
        cleaned = pd.to_numeric(series, errors="coerce").dropna()
        if cleaned.empty:
            return
        fig, ax = plt.subplots(figsize=(8, 4.2))
        ax.hist(cleaned, bins=18, color="#4F8FBF", edgecolor="white", alpha=0.9)
        ax.set_title(title, fontsize=13)
        ax.set_xlabel(series.name or "Value")
        ax.set_ylabel("Count")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        fig.savefig(path, dpi=self.CHART_DPI)
        plt.close(fig)

    def _render_chart(self, path: Path, render: Any) -> bool:
        if path.exists():
            path.unlink()
        render()
        return path.exists()

    def _format_metric_label(self, has_target: bool) -> str:
        return "Outcome Rate" if has_target else "Volume"

    def summarize_dataset(self, dataset_path: str | Path) -> dict[str, Any]:
        adapted = load_analysis_dataset(dataset_path)
        dataset = adapted.dataframe
        schema = self.infer_schema(dataset)
        target_metric = self._target_metric(dataset, schema)
        time_series = self._time_series(dataset, schema)
        primary_series = self._series_by_dimension(dataset, schema, schema.get("primary_dimension"))
        secondary_series = self._series_by_dimension(dataset, schema, schema.get("secondary_dimension"))
        bucket_series, bucket_col = self._bucket_series(dataset, schema)
        missing_share = self._count_missing(dataset)

        target_rate = float(target_metric.mean()) if target_metric is not None else None
        top_time = str(time_series.index[-1]) if time_series is not None and len(time_series) else "No clear time pattern detected"
        top_primary = str(primary_series.index[0]) if primary_series is not None and len(primary_series) else "No dominant segment detected"
        top_secondary = str(secondary_series.index[0]) if secondary_series is not None and len(secondary_series) else "No second grouping detected"
        primary_metric = schema.get("primary_metric")
        metric_summary = None
        if primary_metric and primary_metric in dataset.columns:
            metric_series = pd.to_numeric(dataset[primary_metric], errors="coerce").dropna()
            if not metric_series.empty:
                metric_summary = {
                    "column": primary_metric,
                    "mean": float(metric_series.mean()),
                    "median": float(metric_series.median()),
                }

        top_missing = None
        if not missing_share.empty and float(missing_share.iloc[0]) > 0:
            top_missing = {"column": str(missing_share.index[0]), "missing_pct": round(float(missing_share.iloc[0]), 2)}

        key_findings: list[str] = []
        if target_rate is not None:
            key_findings.append(f"Detected binary outcome rate is {target_rate:.1%}.")
        if schema.get("primary_dimension") and primary_series is not None and len(primary_series):
            key_findings.append(f"Top grouping on {schema['primary_dimension']} is {top_primary}.")
        if schema.get("secondary_dimension") and secondary_series is not None and len(secondary_series):
            key_findings.append(f"Secondary pattern is led by {top_secondary} on {schema['secondary_dimension']}.")
        if schema.get("time") and time_series is not None and len(time_series):
            key_findings.append(f"Time pattern is available through {schema['time']}.")
        if metric_summary is not None:
            key_findings.append(
                f"Primary numeric metric {metric_summary['column']} has mean {metric_summary['mean']:.2f} and median {metric_summary['median']:.2f}."
            )
        if top_missing is not None:
            key_findings.append(f"Largest completeness gap is {top_missing['column']} at {top_missing['missing_pct']:.1f}% missing.")

        return {
            "dataset": dataset,
            "schema": schema,
            "rows": adapted.raw_rows,
            "columns": len(dataset.columns),
            "analysis_rows": adapted.analysis_rows,
            "analysis_grain": adapted.analysis_grain,
            "source_format": adapted.source_format,
            "target_rate": target_rate,
            "top_time": top_time,
            "top_primary_dimension_value": top_primary,
            "top_secondary_dimension_value": top_secondary,
            "primary_metric_summary": metric_summary,
            "bucket_col": bucket_col,
            "time_series": time_series,
            "primary_series": primary_series,
            "secondary_series": secondary_series,
            "bucket_series": bucket_series,
            "top_missing": top_missing,
            "missing_summary": missing_share.head(8),
            "key_findings": key_findings[:6],
        }

    def generate_charts(self, dataset_path: str | Path, chart_prefix: str | None = None) -> dict[str, Path]:
        summary = self.summarize_dataset(dataset_path)
        dataset = summary["dataset"]
        schema = summary["schema"]
        has_target = summary["target_rate"] is not None
        ylabel = self._format_metric_label(has_target)
        prefix = f"{chart_prefix}_" if chart_prefix else ""

        chart_paths: dict[str, Path] = {}
        plt.style.use("seaborn-v0_8-whitegrid")

        missing = summary["missing_summary"]
        if len(missing) and float(missing.iloc[0]) > 0:
            path = self.chart_dir / f"{prefix}missingness.png"
            if self._render_chart(path, lambda: self._plot_bar(missing[missing > 0], "Missingness by Column", path, horizontal=True, ylabel="Missing %")):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["missingness"] = path

        time_series = summary["time_series"]
        if time_series is not None and len(time_series):
            path = self.chart_dir / f"{prefix}time_trend.png"
            if self._render_chart(path, lambda: self._plot_line(time_series, f"Trend Across {schema['time']}", path, ylabel=ylabel)):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["time_trend"] = path

        primary_series = summary["primary_series"]
        if primary_series is not None and len(primary_series):
            path = self.chart_dir / f"{prefix}primary_dimension.png"
            title = f"Top Groups by {schema['primary_dimension']}"
            if self._render_chart(path, lambda: self._plot_bar(primary_series, title, path, ylabel=ylabel)):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["primary_dimension"] = path

        secondary_series = summary["secondary_series"]
        if secondary_series is not None and len(secondary_series):
            path = self.chart_dir / f"{prefix}secondary_dimension.png"
            title = f"Top Groups by {schema['secondary_dimension']}"
            if self._render_chart(path, lambda: self._plot_bar(secondary_series, title, path, horizontal=True, ylabel=ylabel)):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["secondary_dimension"] = path

        bucket_series = summary["bucket_series"]
        if bucket_series is not None and len(bucket_series):
            path = self.chart_dir / f"{prefix}metric_buckets.png"
            title = f"Performance by {summary['bucket_col']} Bucket"
            if self._render_chart(path, lambda: self._plot_line(bucket_series, title, path, ylabel=ylabel)):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["metric_buckets"] = path
        elif schema.get("primary_metric") and schema["primary_metric"] in dataset.columns:
            metric_col = schema["primary_metric"]
            path = self.chart_dir / f"{prefix}metric_distribution.png"
            if self._render_chart(path, lambda: self._plot_hist(dataset[metric_col], f"Distribution of {metric_col}", path)):
                shutil.copy2(path, self.docs_asset_dir / path.name)
                chart_paths["metric_distribution"] = path

        return chart_paths

    def _dataset_title(self, summary: dict[str, Any]) -> str:
        schema = summary["schema"]
        if schema.get("target"):
            return f"{schema['target']} Insight Summary"
        if schema.get("primary_dimension"):
            return f"{schema['primary_dimension']} Performance Review"
        return "Dataset Insight Summary"

    def _recommendations(self, summary: dict[str, Any]) -> list[str]:
        schema = summary["schema"]
        recommendations: list[str] = []
        if schema.get("primary_dimension"):
            recommendations.append(
                f"Prioritize follow-up analysis on the leading {schema['primary_dimension']} group: {summary['top_primary_dimension_value']}."
            )
        if schema.get("secondary_dimension"):
            recommendations.append(
                f"Compare operational or commercial differences across {schema['secondary_dimension']} to confirm whether the pattern is stable."
            )
        if schema.get("primary_metric"):
            recommendations.append(
                f"Track {schema['primary_metric']} as a core KPI because it is the strongest detected numeric driver in this dataset."
            )
        if summary["top_missing"] is not None:
            recommendations.append(
                f"Improve data quality for {summary['top_missing']['column']} before using it in downstream decision-making."
            )
        if not recommendations:
            recommendations.append("Profile additional columns and business context before committing to action based on this dataset alone.")
        recommendations.append("Re-run this presentation after the next data refresh to confirm that the same patterns persist.")
        return recommendations[:4]

    def build_slide_contents(self, dataset_path: str | Path) -> list[SlideContent]:
        summary = self.summarize_dataset(dataset_path)
        schema = summary["schema"]
        metric_summary = summary["primary_metric_summary"]
        target_text = (
            f"Detected binary outcome rate is {summary['target_rate']:.1%}."
            if summary["target_rate"] is not None
            else "No clean binary outcome was detected, so the presentation emphasizes distributions, group structure, and metric patterns."
        )
        metric_text = "No strong numeric metric was detected."
        if metric_summary is not None:
            metric_text = (
                f"{metric_summary['column']} is the lead numeric field with mean {metric_summary['mean']:.2f} "
                f"and median {metric_summary['median']:.2f}."
            )
        quality_text = (
            f"Largest completeness issue is {summary['top_missing']['column']} at {summary['top_missing']['missing_pct']:.1f}% missing."
            if summary["top_missing"] is not None
            else "No major missingness issue stands out in the analyzed dataset."
        )

        return [
            SlideContent(
                title=self._dataset_title(summary),
                bullets=[
                    "This deck was generated from the uploaded dataset without relying on a fixed business template.",
                    (
                        f"The file contains {summary['rows']} raw rows and is summarized at the "
                        f"{summary['analysis_grain']} level across {summary['analysis_rows']} analyzed records."
                    ),
                    target_text,
                ],
                speaker_notes="Open with the dataset scope, the automated analysis grain, and whether the system found a usable binary outcome.",
            ),
            SlideContent(
                title="Detected Structure",
                bullets=[
                    f"Detected time field: {schema.get('time') or 'not confidently identified'}.",
                    f"Primary grouping field: {schema.get('primary_dimension') or 'not confidently identified'}.",
                    f"Secondary grouping field: {schema.get('secondary_dimension') or 'not confidently identified'}.",
                    metric_text,
                ],
                speaker_notes="Explain which columns were inferred as the most useful structural dimensions for the analysis.",
                chart_key="time_trend",
            ),
            SlideContent(
                title="Leading Groups",
                bullets=[
                    f"Leading value for {schema.get('primary_dimension') or 'the primary grouping'} is {summary['top_primary_dimension_value']}.",
                    f"Leading value for {schema.get('secondary_dimension') or 'the secondary grouping'} is {summary['top_secondary_dimension_value']}.",
                    "These groups are the first place to look for concentrated opportunity, demand, or operational differences.",
                    "The best-performing grouping may be worth breaking down further with more business context.",
                ],
                speaker_notes="Use the strongest detected categorical patterns to frame which groups deserve the most attention.",
                chart_key="primary_dimension",
            ),
            SlideContent(
                title="Metric Patterns",
                bullets=[
                    metric_text,
                    (
                        f"Bucket analysis for {summary['bucket_col']} shows how performance changes across the strongest detected metric."
                        if summary["bucket_col"]
                        else "A direct metric distribution is shown because no reliable bucketed performance view was available."
                    ),
                    "This gives a more reusable view of the dataset than assuming a single ecommerce funnel structure.",
                ],
                speaker_notes="Explain how the lead metric behaves and whether higher or lower buckets correspond to stronger outcomes or greater volume.",
                chart_key="metric_buckets" if summary["bucket_col"] else "metric_distribution",
            ),
            SlideContent(
                title="Data Quality",
                bullets=[
                    quality_text,
                    f"Total analyzed columns: {summary['columns']}.",
                    "Data quality affects how confidently we can interpret the strongest patterns in the deck.",
                    "Columns with missing or weak signal should be validated before they drive operational decisions.",
                ],
                speaker_notes="Call out the largest quality risk so the audience understands where the analysis is strongest and where caution is needed.",
                chart_key="missingness",
            ),
            SlideContent(
                title="Recommendations",
                bullets=self._recommendations(summary),
                speaker_notes="Turn the detected structure into a short set of practical next steps grounded in the strongest observed signals.",
                chart_key="secondary_dimension",
            ),
            SlideContent(
                title="Key Takeaways",
                bullets=summary["key_findings"] or [
                    "The dataset contains enough structure to support an automated summary, but the strongest story depends on more business context."
                ],
                speaker_notes="Close with the sharpest findings surfaced automatically from the uploaded dataset.",
            ),
        ]

    def _apply_theme(self, slide) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(251, 252, 250)

    def _add_title(self, slide, title: str) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.8), Inches(0.65))
        paragraph = title_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = title
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(24, 64, 84)

    def _add_bullets(self, slide, bullets: list[str]) -> None:
        text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.18), Inches(5.25), Inches(5.0))
        frame = text_box.text_frame
        frame.word_wrap = True
        for index, bullet in enumerate(bullets[:5]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(7)

    def _add_notes(self, slide, speaker_notes: str) -> None:
        try:
            slide.notes_slide.notes_text_frame.text = speaker_notes
        except Exception:
            footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(8.2), Inches(0.3))
            paragraph = footer.text_frame.paragraphs[0]
            paragraph.text = f"Speaker notes: {speaker_notes}"
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(96, 112, 122)

    def _add_chart(self, slide, chart_path: Path) -> None:
        slide.shapes.add_picture(str(chart_path), Inches(6.05), Inches(1.32), width=Inches(3.18))

    def create_presentation(
        self,
        dataset_path: str | Path,
        output_path: str | Path,
        chart_prefix: str | None = None,
    ) -> tuple[Path, list[SlideContent], dict[str, Path]]:
        chart_paths = self.generate_charts(dataset_path, chart_prefix=chart_prefix)
        slides = self.build_slide_contents(dataset_path)

        presentation = Presentation()
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
        layout = presentation.slide_layouts[6]

        for slide_content in slides:
            slide = presentation.slides.add_slide(layout)
            self._apply_theme(slide)
            self._add_title(slide, slide_content.title)
            wrapped = [textwrap.fill(item, width=57) for item in slide_content.bullets]
            self._add_bullets(slide, wrapped)
            if slide_content.chart_key and slide_content.chart_key in chart_paths:
                self._add_chart(slide, chart_paths[slide_content.chart_key])
            self._add_notes(slide, slide_content.speaker_notes)

        output_path = Path(output_path)
        presentation.save(output_path)
        shutil.copy2(output_path, self.docs_asset_dir / output_path.name)
        return output_path, slides, chart_paths
